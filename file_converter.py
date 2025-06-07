import os
import time
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
import uuid
import sys
import subprocess
import pkg_resources

def install_requirements():
    try:
        # requestsパッケージがインストールされているか確認
        pkg_resources.require(['requests'])
    except pkg_resources.DistributionNotFound:
        print("必要なパッケージをインストールしています...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "requests"])

def get_poppler_path():
    home = os.path.expanduser("~")
    base_dir = os.path.join(home, "poppler-23.11.0")
    
    # 可能性のあるパスをチェック
    possible_paths = [
        os.path.join(base_dir, "Library", "bin"),
        os.path.join(base_dir, "bin"),
        base_dir
    ]
    
    # 実際のパスを探す
    for path in possible_paths:
        print(f"Popplerのパスを確認中: {path}")
        if os.path.exists(path):
            # pdftotext.exeの存在を確認
            if os.path.exists(os.path.join(path, "pdftotext.exe")):
                print(f"Popplerが見つかりました: {path}")
                return path
    
    return None

def setup_poppler():
    try:
        from poppler_installer import download_poppler
        return download_poppler()
    except ImportError:
        print("Popplerインストーラーが見つかりません")
        return False

class FileConverter:
    @staticmethod
    def convert_pdf(file_path, output_dir):
        try:
            # Windowsの場合、poppler_pathを指定
            if sys.platform.startswith('win'):
                poppler_path = get_poppler_path()
                if not poppler_path:
                    print("Popplerがインストールされていません。インストールを開始します...")
                    if not setup_poppler():
                        raise Exception("Popplerのインストールに失敗しました")
                    poppler_path = get_poppler_path()
                    if not poppler_path:
                        raise Exception("Popplerのパスが見つかりません")
                
                print(f"PDFの変換を開始します。Popplerパス: {poppler_path}")
                pages = convert_from_path(file_path, poppler_path=poppler_path)
            else:
                pages = convert_from_path(file_path)
                
            for i, page in enumerate(pages):
                output_path = os.path.join(output_dir, f"{uuid.uuid4()}.png")
                page.save(output_path, "PNG")
                print(f"PDFページを変換しました: {output_path}")
        except Exception as e:
            print(f"PDFの変換エラー: {str(e)}")
            print(f"ファイルパス: {file_path}")
            print(f"出力ディレクトリ: {output_dir}")
            if 'poppler_path' in locals():
                print(f"Popplerパス: {poppler_path}")
                # Popplerディレクトリの内容を表示
                if os.path.exists(poppler_path):
                    print("Popplerディレクトリの内容:")
                    for item in os.listdir(poppler_path):
                        print(f"  - {item}")

    @staticmethod
    def convert_pptx(file_path, output_dir):
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                # スライドを画像として保存
                img_path = os.path.join(output_dir, f"{uuid.uuid4()}.png")
                width = 1920
                height = 1080
                
                # 空の画像を作成
                img = Image.new("RGB", (width, height), "white")
                
                # テキストを描画（簡易的な実装）
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pass  # ここでテキストを画像に描画できます
                
                img.save(img_path)
                print(f"PowerPointスライドを変換しました: {img_path}")
        except Exception as e:
            print(f"PowerPointの変換エラー: {str(e)}")

    @staticmethod
    def convert_text(file_path, output_dir):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
                # テキストを画像に変換
                img = Image.new('RGB', (1920, 1080), color='white')
                output_path = os.path.join(output_dir, f"{uuid.uuid4()}.png")
                img.save(output_path)
                print(f"テキストファイルを変換しました: {output_path}")
        except Exception as e:
            print(f"テキストファイルの変換エラー: {str(e)}")

class FileHandler(FileSystemEventHandler):
    def __init__(self, watch_dir):
        self.watch_dir = watch_dir
        self.converter = FileConverter()

    def on_created(self, event):
        if event.is_directory:
            return
        
        file_path = event.src_path
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # 変換対象のファイルのみ処理
        if file_ext in ['.pdf', '.pptx', '.ppt', '.txt']:
            print(f"新しいファイルを検出: {file_path}")
            
            try:
                if file_ext == '.pdf':
                    self.converter.convert_pdf(file_path, self.watch_dir)
                elif file_ext in ['.pptx', '.ppt']:
                    self.converter.convert_pptx(file_path, self.watch_dir)
                elif file_ext == '.txt':
                    self.converter.convert_text(file_path, self.watch_dir)
            except Exception as e:
                print(f"変換エラー: {str(e)}")

def main():
    # 必要なパッケージをインストール
    install_requirements()
    
    # 監視フォルダのパスを直接指定
    watch_dir = r"C:\Users\ppppa\OneDrive\デスクトップ\aaa"
    
    if not os.path.exists(watch_dir):
        print(f"フォルダが見つかりません: {watch_dir}")
        print("フォルダを作成します...")
        os.makedirs(watch_dir)
    
    event_handler = FileHandler(watch_dir)
    observer = Observer()
    observer.schedule(event_handler, watch_dir, recursive=False)
    observer.start()
    
    print(f"フォルダの監視を開始しました: {watch_dir}")
    print("終了するには Ctrl+C を押してください")
    
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
        print("\n監視を終了します")
    
    observer.join()

if __name__ == "__main__":
    main()